#!/usr/bin/env python3
"""Compile the generated slide images into a 16:9 PowerPoint deck.

Each PNG in images/generated/ becomes one full-bleed slide, in filename order.

Usage:  python3 docs/presentation/build-pptx.py
Output: docs/presentation/claude-code-browser.ro.pptx
"""
import glob
import os

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(HERE, "images", "generated")
OUT = os.path.join(HERE, "claude-code-browser.ro.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]       # fully blank layout

images = sorted(glob.glob(os.path.join(IMG_DIR, "slide-*.png")))
if not images:
    raise SystemExit(f"No slides found in {IMG_DIR}")

for img in images:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print(f"Wrote {len(images)} slides to {os.path.relpath(OUT, os.getcwd())}")
for i, img in enumerate(images):
    print(f"  {i}: {os.path.basename(img)}")
